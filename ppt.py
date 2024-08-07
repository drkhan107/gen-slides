import os
from typing import List, Dict, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pptx.shapes.base import BaseShape

class SlideDeck:
    """
    A class to create and manage PowerPoint presentations.
    """

    def __init__(self, output_folder: str = "generated"):
        """
        Initialize the SlideDeck class.

        Args:
            output_folder (str): The folder where the presentation will be saved.
        """
        self.prs = Presentation()
        self.output_folder = output_folder

    def add_slide(self, slide_data: Dict[str, Union[str, List[str], List[List[str]]]]) -> None:
        """
        Add a new slide to the presentation.

        Args:
            slide_data (Dict): A dictionary containing the slide content.
        """
        prs = self.prs
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        # Add title
        title_shape = shapes.title
        title_shape.text = slide_data.get("title_text", "")

        # Add body text
        if "text" in slide_data:
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            for bullet in slide_data.get("text", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

                if "p1" in slide_data:
                    p = tf.add_paragraph()
                    p.text = slide_data.get("p1", "")
                    p.level = 1

        # Add images
        if "img_path" in slide_data:
            cur_left = 6
            for img_path in slide_data.get("img_path", []):
                try:
                    top = Inches(2)
                    left = Inches(cur_left)
                    height = Inches(4)
                    pic = slide.shapes.add_picture(img_path, left, top, height=height)
                    cur_left += 1
                except FileNotFoundError:
                    print(f"Warning: Image file not found: {img_path}")

        # Add table
        if "table" in slide_data:
            self.add_table(slide, slide_data["table"])

    def add_title_slide(self, title_page_data: Dict[str, str]) -> None:
        """
        Add a title slide to the presentation.

        Args:
            title_page_data (Dict): A dictionary containing the title slide content.
        """
        prs = self.prs
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if "title_text" in title_page_data:
            title.text = title_page_data.get("title_text")
        if "subtitle_text" in title_page_data:
            subtitle.text = title_page_data.get("subtitle_text")

    def add_table(self, slide, table_data: List[List[str]]) -> None:
        """
        Add a table to the given slide.

        Args:
            slide (Slide): The slide to add the table to.
            table_data (List[List[str]]): The data for the table.
        """
        # Determine the maximum number of columns
        max_cols = max(len(row) for row in table_data)
        rows = len(table_data)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)

        table = slide.shapes.add_table(rows, max_cols, left, top, width, height).table

        # Set column widths
        first_col_width = 3.5
        remaining_width = 9 - first_col_width
        other_col_width = remaining_width / (max_cols - 1) if max_cols > 1 else remaining_width
        
        table.columns[0].width = Inches(first_col_width)
        for i in range(1, max_cols):
            table.columns[i].width = Inches(other_col_width)

        # Populate the table with data
        for i, row in enumerate(table_data):
            if len(row) < max_cols:
                row.insert(0, " ")
            for j, cell in enumerate(row):
                table.cell(i, j).text = str(cell)
                paragraph = table.cell(i, j).text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 or i == 0 else PP_ALIGN.LEFT

        # Style the header row
        for cell in table.rows[0].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue color
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
            cell.text_frame.paragraphs[0].font.bold = True

        # Style the first column
        for i in range(1, rows):
            cell = table.cell(i, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray
            cell.text_frame.paragraphs[0].font.bold = True

    def create_presentation(self, title_slide_info: Dict[str, str], slide_pages_data: List[Dict[str, Union[str, List[str], List[List[str]]]]] = []) -> str:
        """
        Create a complete presentation.

        Args:
            title_slide_info (Dict): Information for the title slide.
            slide_pages_data (List[Dict]): Data for all other slides.

        Returns:
            str: The file path of the saved presentation.

        Raises:
            OSError: If there's an error creating or saving the file.
            ValueError: If the input data is invalid.
        """
        try:
            # Generate file name from title
            file_name = title_slide_info.get("title_text", "presentation").\
                lower().replace(",", "").replace(":", "").replace(" ", "-")
            file_name += ".pptx"
            file_name = os.path.join(self.output_folder, file_name)

            # Create output folder if it doesn't exist
            os.makedirs(self.output_folder, exist_ok=True)

            # Add title slide
            self.add_title_slide(title_slide_info)

            # Add content slides
            for slide_data in slide_pages_data:
                self.add_slide(slide_data)

            # Save the presentation
            self.prs.save(file_name)
            return file_name
        except OSError as e:
            raise OSError(f"Error creating or saving the presentation: {str(e)}")
        except ValueError as e:
            raise ValueError(f"Invalid input data: {str(e)}")
        except Exception as e:
            raise Exception(f"An unexpected error occurred: {str(e)}")